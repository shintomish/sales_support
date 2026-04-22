"""
scripts/output/200_Flowmap.pptx 生成スクリプト

sales_support システムの業務フローマップを PowerPoint で作成する。
出力先: scripts/output/200_Flowmap.pptx (このスクリプトと同じ階層の output/)

使い方:
    cd ~/sales_support
    python3 scripts/gen_flowmap.py

出力ディレクトリが存在しない場合は自動作成される。
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── カラーパレット ──────────────────────────────────────────
C_BLUE_DARK  = RGBColor(0x1F, 0x45, 0x7C)   # 濃紺（タイトル帯）
C_BLUE       = RGBColor(0x2E, 0x75, 0xB6)   # 青（実装済みノード）
C_GREEN      = RGBColor(0x37, 0x86, 0x4A)   # 緑（完了・送信）
C_ORANGE     = RGBColor(0xE3, 0x7B, 0x1F)   # オレンジ（未接続・課題）
C_GRAY       = RGBColor(0x76, 0x76, 0x76)   # グレー（補足）
C_LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)   # 薄青（背景帯）
C_LIGHT_GREEN= RGBColor(0xC6, 0xEF, 0xCE)   # 薄緑（背景帯）
C_LIGHT_ORG  = RGBColor(0xFC, 0xE4, 0xD6)   # 薄オレンジ（背景帯）
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def box(slide, x, y, w, h,
        text="", font_size=None, bold=False,
        fill=C_BLUE, text_color=C_WHITE,
        border_color=None, border_width=Pt(1),
        align=PP_ALIGN.CENTER, valign=None,
        shadow=False, radius=None):
    """汎用テキストボックス（角丸対応）"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    # 塗り
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    # 枠線
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # テキスト
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size if font_size is not None else Pt(11)
    run.font.bold = bold
    run.font.color.rgb = text_color
    # 上下中央
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def arrow(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5), dashed=False):
    """直線矢印"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    if dashed:
        connector.line.dash_style = 4  # DASH
    return connector


def title_bar(slide, title, subtitle=""):
    """スライドタイトル帯"""
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.7),
        text="", fill=C_BLUE_DARK)
    box(slide, Inches(0.2), Inches(0.05), Inches(10), Inches(0.6),
        text=title, font_size=Pt(20), bold=True,
        fill=C_BLUE_DARK, text_color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, Inches(10.2), Inches(0.1), Inches(2.9), Inches(0.5),
            text=subtitle, font_size=Pt(9),
            fill=C_BLUE_DARK, text_color=RGBColor(0xBB,0xCC,0xEE),
            align=PP_ALIGN.RIGHT)


def label(slide, x, y, w, h, text, font_size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER):
    """透明背景ラベル"""
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    return shape


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド1: 表紙
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_cover(prs):
    sl = add_slide(prs)
    box(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
        fill=C_BLUE_DARK)
    box(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
        text="業務フローマップ", font_size=Pt(36), bold=True,
        fill=C_BLUE_DARK, text_color=C_WHITE)
    box(sl, Inches(1.5), Inches(3.2), Inches(10), Inches(0.7),
        text="sales_support — 案件メール・技術者管理・マッチング・配信の連携全体図",
        font_size=Pt(16),
        fill=C_BLUE_DARK, text_color=RGBColor(0xBB,0xCC,0xEE))
    box(sl, Inches(1.5), Inches(4.2), Inches(4), Inches(0.5),
        text="2026-04-11",
        font_size=Pt(13),
        fill=C_BLUE_DARK, text_color=RGBColor(0x88,0xAA,0xCC))

    # 凡例
    box(sl, Inches(1.5), Inches(5.2), Inches(1.8), Inches(0.4),
        text="✅ 実装済み", font_size=Pt(11), fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(3.5), Inches(5.2), Inches(1.8), Inches(0.4),
        text="🔄 SES審査中", font_size=Pt(11), fill=C_GREEN, text_color=C_WHITE)
    box(sl, Inches(5.5), Inches(5.2), Inches(2.0), Inches(0.4),
        text="⚠ 未接続（今後の課題）", font_size=Pt(11),
        fill=C_ORANGE, text_color=C_WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド2: システム全体像
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_overview(prs):
    sl = add_slide(prs)
    title_bar(sl, "システム全体像 — 5つの機能エリアと連携状況", "Overview")

    # 背景帯: 受信
    box(sl, Inches(0.2), Inches(0.85), Inches(2.6), Inches(6.4),
        fill=C_LIGHT_BLUE, border_color=C_BLUE, border_width=Pt(0.75))
    label(sl, Inches(0.3), Inches(0.88), Inches(2.4), Inches(0.3),
          "① 受信・分類", font_size=Pt(9), color=C_BLUE_DARK)

    # 背景帯: 案件系
    box(sl, Inches(3.1), Inches(0.85), Inches(3.0), Inches(3.0),
        fill=C_LIGHT_BLUE, border_color=C_BLUE, border_width=Pt(0.75))
    label(sl, Inches(3.2), Inches(0.88), Inches(2.8), Inches(0.3),
          "② 案件メール管理", font_size=Pt(9), color=C_BLUE_DARK)

    # 背景帯: 技術者系
    box(sl, Inches(3.1), Inches(4.1), Inches(3.0), Inches(2.9),
        fill=C_LIGHT_BLUE, border_color=C_BLUE, border_width=Pt(0.75))
    label(sl, Inches(3.2), Inches(4.13), Inches(2.8), Inches(0.3),
          "③ 技術者メール管理", font_size=Pt(9), color=C_BLUE_DARK)

    # 背景帯: マスタ・マーケット
    box(sl, Inches(6.4), Inches(0.85), Inches(3.2), Inches(6.4),
        fill=C_LIGHT_GREEN, border_color=C_GREEN, border_width=Pt(0.75))
    label(sl, Inches(6.5), Inches(0.88), Inches(3.0), Inches(0.3),
          "④ マスタ管理 / マーケット", font_size=Pt(9), color=C_GREEN)

    # 背景帯: 送信
    box(sl, Inches(9.9), Inches(0.85), Inches(3.2), Inches(6.4),
        fill=C_LIGHT_ORG, border_color=C_ORANGE, border_width=Pt(0.75))
    label(sl, Inches(10.0), Inches(0.88), Inches(3.0), Inches(0.3),
          "⑤ 送信チャネル", font_size=Pt(9), color=C_ORANGE)

    # ── ① 受信エリア ──────────────────────────────
    box(sl, Inches(0.3), Inches(1.3), Inches(2.2), Inches(0.55),
        text="Gmail受信", font_size=Pt(11), bold=True,
        fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(0.3), Inches(2.2), Inches(2.2), Inches(0.55),
        text="案件メール\n(自動分類・AI解析)",
        font_size=Pt(10), fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(0.3), Inches(4.2), Inches(2.2), Inches(0.55),
        text="技術者メール\n(自動分類・AI解析)",
        font_size=Pt(10), fill=C_BLUE, text_color=C_WHITE)

    arrow(sl, Inches(1.4), Inches(1.85), Inches(1.4), Inches(2.2),
          color=C_BLUE, width=Pt(1.5))
    arrow(sl, Inches(1.4), Inches(1.85), Inches(1.4), Inches(4.2),
          color=C_BLUE, width=Pt(1.5))
    # 分岐点
    box(sl, Inches(0.3), Inches(1.85), Inches(2.2), Inches(0.35),
        text="分類エンジン", font_size=Pt(9),
        fill=RGBColor(0x5B,0xA3,0xD6), text_color=C_WHITE)

    # ── ① → ② 矢印 ──
    arrow(sl, Inches(2.5), Inches(2.48), Inches(3.1), Inches(2.48),
          color=C_BLUE, width=Pt(2))
    arrow(sl, Inches(2.5), Inches(4.48), Inches(3.1), Inches(4.48),
          color=C_BLUE, width=Pt(2))

    # ── ② 案件メール ──────────────────────────────
    box(sl, Inches(3.2), Inches(1.3), Inches(2.8), Inches(0.55),
        text="AIスコアリング", font_size=Pt(10),
        fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(3.2), Inches(2.1), Inches(2.8), Inches(0.55),
        text="案件メール一覧\n確認・編集・優先度付け",
        font_size=Pt(10), fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(3.2), Inches(2.9), Inches(2.8), Inches(0.55),
        text="マッチング\n(自社技術者との照合)",
        font_size=Pt(10), fill=C_BLUE, text_color=C_WHITE)

    arrow(sl, Inches(4.6), Inches(2.48), Inches(4.6), Inches(2.1),
          color=C_BLUE, width=Pt(1.5))
    arrow(sl, Inches(4.6), Inches(2.65), Inches(4.6), Inches(2.9),
          color=C_BLUE, width=Pt(1.5))

    # ── ③ 技術者メール ─────────────────────────────
    box(sl, Inches(3.2), Inches(4.2), Inches(2.8), Inches(0.55),
        text="AIスコアリング", font_size=Pt(10),
        fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(3.2), Inches(5.0), Inches(2.8), Inches(0.55),
        text="技術者メール一覧\n確認・ステータス管理",
        font_size=Pt(10), fill=C_BLUE, text_color=C_WHITE)
    box(sl, Inches(3.2), Inches(5.8), Inches(2.8), Inches(0.55),
        text="→ 技術者登録（手動）\n⚠ 自動連携なし",
        font_size=Pt(9), fill=C_ORANGE, text_color=C_WHITE)

    arrow(sl, Inches(4.6), Inches(4.75), Inches(4.6), Inches(5.0),
          color=C_BLUE, width=Pt(1.5))
    arrow(sl, Inches(4.6), Inches(5.55), Inches(4.6), Inches(5.8),
          color=C_ORANGE, width=Pt(1.5), dashed=True)

    # ── ④ マスタ ──────────────────────────────────
    box(sl, Inches(6.5), Inches(1.3), Inches(2.8), Inches(0.6),
        text="技術者管理\n(Engineer マスタ)",
        font_size=Pt(10), bold=True, fill=C_GREEN, text_color=C_WHITE)
    box(sl, Inches(6.5), Inches(2.2), Inches(2.8), Inches(0.6),
        text="案件マーケット\n(PublicProject)",
        font_size=Pt(10), bold=True, fill=C_GREEN, text_color=C_WHITE)
    box(sl, Inches(6.5), Inches(3.1), Inches(2.8), Inches(0.6),
        text="マッチングスコア\n(MatchingService)",
        font_size=Pt(10), fill=C_GREEN, text_color=C_WHITE)
    box(sl, Inches(6.5), Inches(4.2), Inches(2.8), Inches(0.6),
        text="配信先リスト\n(DeliveryAddress)",
        font_size=Pt(10), fill=C_GREEN, text_color=C_WHITE)

    # マッチング矢印
    arrow(sl, Inches(6.0), Inches(3.4), Inches(6.5), Inches(3.4),
          color=C_GREEN, width=Pt(2))
    arrow(sl, Inches(7.9), Inches(2.5), Inches(7.9), Inches(3.1),
          color=C_GREEN, width=Pt(1.5))
    arrow(sl, Inches(7.9), Inches(1.6), Inches(7.9), Inches(2.2),
          color=C_ORANGE, width=Pt(1.5), dashed=True)

    # 未接続注記
    label(sl, Inches(6.5), Inches(3.72), Inches(2.8), Inches(0.3),
          "⚠ マーケット→送信 未接続", font_size=Pt(8), color=C_ORANGE)

    # ── ⑤ 送信 ───────────────────────────────────
    box(sl, Inches(10.0), Inches(1.3), Inches(2.8), Inches(0.6),
        text="提案メール送信\n(個別 / 一括)",
        font_size=Pt(10), bold=True, fill=C_GREEN, text_color=C_WHITE)
    box(sl, Inches(10.0), Inches(2.4), Inches(2.8), Inches(0.6),
        text="配信キャンペーン\n(DeliveryCampaign)",
        font_size=Pt(10), bold=True,
        fill=RGBColor(0xE3,0x9B,0x1F), text_color=C_WHITE)
    label(sl, Inches(10.0), Inches(3.05), Inches(2.8), Inches(0.3),
          "🔄 Amazon SES審査中", font_size=Pt(9),
          color=RGBColor(0xB0,0x60,0x00))

    # マッチング → 提案送信
    arrow(sl, Inches(9.3), Inches(3.4), Inches(10.0), Inches(1.6),
          color=C_GREEN, width=Pt(2))
    # 配信先 → キャンペーン
    arrow(sl, Inches(9.3), Inches(4.5), Inches(10.0), Inches(2.7),
          color=C_GREEN, width=Pt(2))
    # 案件メール → キャンペーン
    arrow(sl, Inches(6.0), Inches(2.5), Inches(10.0), Inches(2.5),
          color=RGBColor(0xE3,0x9B,0x1F), width=Pt(1.5), dashed=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド3: フローA — 外部案件 → 自社技術者提案
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_flow_a(prs):
    sl = add_slide(prs)
    title_bar(sl, "フローA: 外部案件メール → 自社技術者の提案", "✅ 実装済み")

    # ステップボックス
    steps = [
        ("①\nパートナー企業から\n案件紹介メール受信", C_BLUE),
        ("②\nGmail同期\n自動分類", C_BLUE),
        ("③\nAIスコアリング\n(重要度・適合度)", C_BLUE),
        ("④\n営業が確認\n手動修正・優先度付け", C_BLUE),
        ("⑤\nマッチング\n自社技術者を照合", C_BLUE),
        ("⑥\n提案メール生成\n(Claude AIで文面作成)", C_BLUE),
        ("⑦\n送信\n(個別 / 一括)", C_GREEN),
    ]

    x = Inches(0.4)
    bw = Inches(1.6)
    bh = Inches(1.4)
    by = Inches(2.5)
    gap = Inches(0.25)

    for i, (txt, color) in enumerate(steps):
        bx = x + (bw + gap) * i
        box(sl, bx, by, bw, bh, text=txt,
            font_size=Pt(10), fill=color, text_color=C_WHITE)
        if i < len(steps) - 1:
            arrow(sl,
                  bx + bw, by + bh/2,
                  bx + bw + gap, by + bh/2,
                  color=C_BLUE, width=Pt(2))

    # データモデル注記
    models = [
        "Email\n(受信メール)",
        "ProjectMail\nSource",
        "ProjectMail\nSource",
        "ProjectMail\nSource\n+ status",
        "Matching\nService",
        "Claude\nAPI",
        "MailSend\nHistory",
    ]
    note_y = Inches(4.2)
    for i, m in enumerate(models):
        bx = x + (bw + gap) * i
        box(sl, bx, note_y, bw, Inches(0.8),
            text=m, font_size=Pt(8),
            fill=C_LIGHT_BLUE, text_color=C_BLUE_DARK,
            border_color=C_BLUE, border_width=Pt(0.5))
        arrow(sl, bx + bw/2, by + bh, bx + bw/2, note_y,
              color=C_GRAY, width=Pt(1))

    # ステータス遷移
    box(sl, Inches(0.4), Inches(5.5), Inches(12.0), Inches(0.6),
        text="ステータス遷移:  new → review → proposed → interview → won / lost",
        font_size=Pt(11), fill=C_LIGHT_BLUE, text_color=C_BLUE_DARK,
        border_color=C_BLUE, border_width=Pt(0.5))

    # タイトル補足
    label(sl, Inches(0.4), Inches(1.1), Inches(12), Inches(0.4),
          "パートナー企業から届いた案件紹介メールを AI でスコアリングし、自社の技術者とマッチングして提案返信するメインフロー",
          font_size=Pt(11), color=C_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド4: フローB — 外部技術者メール（課題あり）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_flow_b(prs):
    sl = add_slide(prs)
    title_bar(sl, "フローB: 外部技術者メール → 自社案件へのアサイン", "⚠ 一部未実装")

    label(sl, Inches(0.4), Inches(1.0), Inches(12), Inches(0.4),
          "外部BP企業から届いた技術者紹介メールを確認し、案件にアサインするフロー。受信〜スコアリングは実装済みだが、後半の連携が未実装。",
          font_size=Pt(11), color=C_GRAY)

    steps_done = [
        ("①\n技術者紹介メール\n受信", C_BLUE),
        ("②\nGmail同期\n自動分類", C_BLUE),
        ("③\nAIスコアリング\n(スキル・稼働日)", C_BLUE),
        ("④\n技術者メール一覧\nで確認・編集", C_BLUE),
    ]
    steps_todo = [
        ("⑤\n自社案件と\nマッチング", C_ORANGE),
        ("⑥\n技術者管理へ\n自動登録", C_ORANGE),
        ("⑦\n案件アサイン\n・提案", C_ORANGE),
    ]

    x = Inches(0.4)
    bw = Inches(1.55)
    bh = Inches(1.4)
    by = Inches(2.3)
    gap = Inches(0.22)

    all_steps = steps_done + steps_todo
    for i, (txt, color) in enumerate(all_steps):
        bx = x + (bw + gap) * i
        box(sl, bx, by, bw, bh, text=txt,
            font_size=Pt(10), fill=color, text_color=C_WHITE)
        if i < len(all_steps) - 1:
            c = C_BLUE if i < len(steps_done) - 1 else C_ORANGE
            arrow(sl,
                  bx + bw, by + bh/2,
                  bx + bw + gap, by + bh/2,
                  color=c, width=Pt(2),
                  dashed=(i >= len(steps_done) - 1))

    # 実装済み / 未実装 の区切り線
    cut_x = x + (bw + gap) * len(steps_done) - gap/2
    arrow(sl, cut_x, Inches(2.0), cut_x, Inches(4.1),
          color=C_ORANGE, width=Pt(2))
    label(sl, cut_x - Inches(0.5), Inches(1.9), Inches(1.5), Inches(0.3),
          "← 実装済み ｜ 未実装 →", font_size=Pt(9), color=C_ORANGE)

    # 課題ボックス
    box(sl, Inches(0.4), Inches(4.8), Inches(12.2), Inches(1.6),
        text="", fill=C_LIGHT_ORG,
        border_color=C_ORANGE, border_width=Pt(0.75))
    label(sl, Inches(0.6), Inches(4.85), Inches(11.8), Inches(0.3),
          "▼ 未実装による課題", font_size=Pt(11), color=C_ORANGE)
    issues = [
        "• EngineerMailSource → Engineer マスタへの自動登録ボタンがない（手動で別途入力が必要）",
        "• 技術者メールを見ながら「自社のどの案件に合うか」を確認する画面がない",
        "• ステータスを registered に変更しても、Engineerレコードは作成されない（DBレベルで未接続）",
    ]
    for j, iss in enumerate(issues):
        label(sl, Inches(0.8), Inches(5.2) + Inches(0.32)*j,
              Inches(11.5), Inches(0.3),
              iss, font_size=Pt(10), color=C_BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド5: フローC — 案件マーケット
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_flow_c(prs):
    sl = add_slide(prs)
    title_bar(sl, "フローC: 案件マーケット → 技術者レコメンド", "⚠ 送信導線なし")

    label(sl, Inches(0.4), Inches(1.0), Inches(12), Inches(0.4),
          "自社が抱える案件を PublicProject として登録し、合う技術者をAIでレコメンドするフロー。レコメンドまでは動作するが送信導線がない。",
          font_size=Pt(11), color=C_GRAY)

    steps = [
        ("①\nPublicProject\n登録", C_GREEN, True),
        ("②\n必須スキル\nを設定", C_GREEN, True),
        ("③\nMatchingService\nでスコア計算", C_GREEN, True),
        ("④\nおすすめ技術者\n一覧を確認", C_GREEN, True),
        ("⑤\n提案メール\nを送信", C_ORANGE, False),
    ]

    x = Inches(1.0)
    bw = Inches(2.0)
    bh = Inches(1.5)
    by = Inches(2.2)
    gap = Inches(0.5)

    for i, (txt, color, done) in enumerate(steps):
        bx = x + (bw + gap) * i
        box(sl, bx, by, bw, bh, text=txt,
            font_size=Pt(11), fill=color, text_color=C_WHITE)
        if not done:
            label(sl, bx, by + bh + Inches(0.05), bw, Inches(0.3),
                  "⚠ 未実装", font_size=Pt(10), color=C_ORANGE)
        if i < len(steps) - 1:
            arrow(sl, bx + bw, by + bh/2, bx + bw + gap, by + bh/2,
                  color=C_GREEN if done else C_ORANGE,
                  width=Pt(2), dashed=(not done))

    # スコア計算の詳細
    box(sl, Inches(1.0), Inches(4.5), Inches(11.0), Inches(1.8),
        text="", fill=C_LIGHT_GREEN,
        border_color=C_GREEN, border_width=Pt(0.75))
    label(sl, Inches(1.2), Inches(4.55), Inches(10.5), Inches(0.3),
          "▼ MatchingService スコア計算ロジック（重みは仮値 → 実運用後チューニング推奨）",
          font_size=Pt(10), color=C_GREEN)
    score_items = [
        "スキル適合度 50%  — 必須スキルの経験年数 × 重要度で計算",
        "単価適合度   25%  — 案件単価と技術者希望単価の重複範囲で計算",
        "勤務地適合度 15%  — remote/hybrid/office の組み合わせで計算",
        "稼働時期適合度 10% — 案件開始日 vs 技術者稼働可能日の差分で計算",
    ]
    for j, item in enumerate(score_items):
        label(sl, Inches(1.4), Inches(4.9) + Inches(0.28)*j,
              Inches(10.5), Inches(0.28),
              item, font_size=Pt(10), color=C_BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド6: フローD — 一括配信
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_flow_d(prs):
    sl = add_slide(prs)
    title_bar(sl, "フローD: 配信先リスト → 一括メール配信", "🔄 Amazon SES審査中")

    label(sl, Inches(0.4), Inches(1.0), Inches(12), Inches(0.4),
          "登録済み配信先リストに対して案件情報を一括配信するフロー。実装は完了しており Amazon SES の本番審査通過後に稼働開始。",
          font_size=Pt(11), color=C_GRAY)

    steps = [
        ("①\n配信先リスト\n管理\n(DeliveryAddress)",
         RGBColor(0x37,0x86,0x4A), True),
        ("②\n配信対象の\n案件メールを選択\n(ProjectMailSource)",
         RGBColor(0x37,0x86,0x4A), True),
        ("③\n件名・本文\nを入力\n(テンプレート使用可)",
         RGBColor(0x37,0x86,0x4A), True),
        ("④\nDeliveryCampaign\nService\n一括送信実行",
         RGBColor(0xE3,0x7B,0x1F), True),
        ("⑤\n送信履歴\n成功/失敗\nを記録",
         RGBColor(0x37,0x86,0x4A), True),
    ]

    x = Inches(0.5)
    bw = Inches(2.1)
    bh = Inches(1.8)
    by = Inches(2.0)
    gap = Inches(0.3)

    for i, (txt, color, _) in enumerate(steps):
        bx = x + (bw + gap) * i
        box(sl, bx, by, bw, bh, text=txt,
            font_size=Pt(10), fill=color, text_color=C_WHITE)
        if i < len(steps) - 1:
            arrow(sl, bx + bw, by + bh/2, bx + bw + gap, by + bh/2,
                  color=color, width=Pt(2))

    # SES審査状況
    box(sl, Inches(0.5), Inches(4.3), Inches(12.2), Inches(0.7),
        text="🔄  Amazon SES 本番送信申請中（2026-04-10 申請 / 承認まで1〜3営業日）  →  承認後すぐ稼働開始予定",
        font_size=Pt(12), bold=True,
        fill=RGBColor(0xFF,0xF0,0xD0),
        border_color=RGBColor(0xE3,0x7B,0x1F),
        border_width=Pt(1.5),
        text_color=RGBColor(0x80,0x40,0x00))

    # 注意点
    box(sl, Inches(0.5), Inches(5.2), Inches(12.2), Inches(1.8),
        text="", fill=C_LIGHT_ORG,
        border_color=C_ORANGE, border_width=Pt(0.75))
    label(sl, Inches(0.7), Inches(5.25), Inches(11.5), Inches(0.3),
          "▼ 審査通過後に確認すべき事項", font_size=Pt(10), color=C_ORANGE)
    notes = [
        "• バウンス・苦情のハンドリング: 送信失敗アドレスを is_active=false に自動更新する仕組みが必要",
        "• 重複送信防止: 同一キャンペーンの再送ガードがあるか確認",
        "• 送信レート制限: SES の sending quota (初期は低め) に合わせてバッチ分割を検討",
    ]
    for j, n in enumerate(notes):
        label(sl, Inches(0.9), Inches(5.6) + Inches(0.35)*j,
              Inches(11.5), Inches(0.3),
              n, font_size=Pt(10), color=C_BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド7: 課題と改善ロードマップ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_issues(prs):
    sl = add_slide(prs)
    title_bar(sl, "断絶箇所と改善ロードマップ", "Next Steps")

    label(sl, Inches(0.4), Inches(0.9), Inches(12), Inches(0.35),
          "現状の未接続ポイントを優先度順に整理し、実装コストと業務インパクトで判断する",
          font_size=Pt(11), color=C_GRAY)

    issues = [
        {
            "prio": "P1",
            "title": "技術者メール → 技術者管理への「一発登録」ボタン",
            "impact": "高",
            "cost": "低",
            "detail": "EngineerMailSource の内容から Engineer レコードを自動生成する API アクションを追加。\n添付スキルシートがあれば parseSkillSheet と組み合わせる。",
            "color": RGBColor(0xC0,0x00,0x00),
        },
        {
            "prio": "P2",
            "title": "技術者メール上での自社案件マッチング表示",
            "impact": "高",
            "cost": "中",
            "detail": "「この技術者に合う自社案件」を ProjectMailSource と同じ仕組みでスコアリングして表示。\nEngineerMailMatchingService として実装。",
            "color": C_ORANGE,
        },
        {
            "prio": "P3",
            "title": "案件マーケット → 提案メール送信の導線",
            "impact": "中",
            "cost": "低",
            "detail": "MatchingController のレコメンド結果から直接 sendProposal を呼び出せる UI / API を追加。\n既存の ProposalMail を流用するだけで実装可能。",
            "color": C_BLUE,
        },
        {
            "prio": "P4",
            "title": "配信管理のバウンス自動処理",
            "impact": "中",
            "cost": "中",
            "detail": "Amazon SES の SNS 通知を受け取り、バウンス・苦情発生時に DeliveryAddress.is_active を自動で false にする。\nSES 審査通過後に着手。",
            "color": C_GRAY,
        },
    ]

    for i, iss in enumerate(issues):
        by = Inches(1.5) + Inches(1.4) * i
        # 優先度バッジ
        box(sl, Inches(0.3), by, Inches(0.55), Inches(0.55),
            text=iss["prio"], font_size=Pt(11), bold=True,
            fill=iss["color"], text_color=C_WHITE)
        # タイトル
        box(sl, Inches(0.95), by, Inches(7.5), Inches(0.55),
            text=iss["title"], font_size=Pt(11), bold=True,
            fill=RGBColor(0xF2,0xF2,0xF2), text_color=C_BLACK,
            border_color=iss["color"], border_width=Pt(1),
            align=PP_ALIGN.LEFT)
        # インパクト・コスト
        box(sl, Inches(8.55), by, Inches(1.0), Inches(0.55),
            text=f"業務\nインパクト: {iss['impact']}",
            font_size=Pt(8), fill=C_LIGHT_GREEN, text_color=C_BLACK,
            border_color=C_GREEN, border_width=Pt(0.5))
        box(sl, Inches(9.65), by, Inches(1.0), Inches(0.55),
            text=f"実装\nコスト: {iss['cost']}",
            font_size=Pt(8), fill=C_LIGHT_BLUE, text_color=C_BLACK,
            border_color=C_BLUE, border_width=Pt(0.5))
        # 詳細
        box(sl, Inches(0.95), by + Inches(0.58), Inches(9.7), Inches(0.72),
            text=iss["detail"], font_size=Pt(9),
            fill=C_WHITE, text_color=C_BLACK,
            border_color=RGBColor(0xCC,0xCC,0xCC), border_width=Pt(0.5),
            align=PP_ALIGN.LEFT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_flow_a(prs)
    slide_flow_b(prs)
    slide_flow_c(prs)
    slide_flow_d(prs)
    slide_issues(prs)

    # スクリプト自身の位置から output/ を解決(ポータブル)
    output_dir = Path(__file__).resolve().parent / "output"
    output_dir.mkdir(parents=True, exist_ok=True)
    out = output_dir / "200_Flowmap.pptx"

    prs.save(str(out))
    print(f"✅ 保存完了: {out}")


if __name__ == "__main__":
    main()
